"""Build a simple, non-technical client presentation for Insight Healthcare Chatbot."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ---------- Brand palette (calm, healthcare-friendly) ----------
NAVY    = RGBColor(0x0F, 0x2A, 0x4A)   # primary dark
TEAL    = RGBColor(0x1F, 0x8F, 0x8F)   # accent
SKY     = RGBColor(0xE6, 0xF2, 0xF7)   # background tint
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SLATE   = RGBColor(0x4A, 0x5A, 0x6A)   # body text
SOFT    = RGBColor(0xD9, 0xE6, 0xEE)   # card border
AMBER   = RGBColor(0xE8, 0x9B, 0x3C)   # warm highlight
MINT    = RGBColor(0xC7, 0xE9, 0xD9)   # soft pill bg
ROSE    = RGBColor(0xF7, 0xD7, 0xD7)   # refuse pill bg


# ---------- Helpers ----------
def set_slide_bg(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height,
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height,
             fill=WHITE, border=SOFT, shadow=True, radius=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    if not shadow:
        shape.shadow.inherit = False
    return shape


def add_pill(slide, left, top, width, height, text,
             fill=MINT, text_color=NAVY, size=11, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return pill


def add_flow_box(slide, left, top, width, height, title,
                 subtitle=None, fill=WHITE, border=TEAL,
                 title_color=NAVY, subtitle_color=SLATE,
                 title_size=14, subtitle_size=10):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.adjustments[0] = 0.15
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Calibri"
        r2.font.size = Pt(subtitle_size)
        r2.font.color.rgb = subtitle_color
    return box


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    # arrowhead via XML
    from pptx.oxml.ns import qn
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    else:
        tail.set('type', 'triangle')
    return line


def slide_header(slide, eyebrow, title, subtitle=None):
    # top accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, prs.slide_width, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    add_text(slide, Inches(0.55), Inches(0.32), Inches(10), Inches(0.3),
             eyebrow.upper(), size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=SLATE)

    # footer
    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             "Insight Healthcare  •  AI Front Desk Assistant",
             size=9, color=SLATE)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.7), Inches(0.3),
             "clinic.callsphere.site", size=9, color=TEAL, align=PP_ALIGN.RIGHT)


# ---------- Build presentation ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# =========================================================
# SLIDE 1 — Title
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, NAVY)

# decorative blob
blob = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(9.5), Inches(-2), Inches(7), Inches(7))
blob.fill.solid(); blob.fill.fore_color.rgb = TEAL
blob.line.fill.background()

blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(-2), Inches(5), Inches(5), Inches(5))
blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x60)
blob2.line.fill.background()

add_text(s, Inches(0.8), Inches(2.0), Inches(10), Inches(0.5),
         "INSIGHT HEALTHCARE", size=14, bold=True, color=TEAL)
add_text(s, Inches(0.8), Inches(2.55), Inches(11), Inches(1.4),
         "Your Clinic's\nAI Front Desk", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.6), Inches(10), Inches(0.6),
         "A simple chatbot that answers patient questions —",
         size=20, color=WHITE)
add_text(s, Inches(0.8), Inches(4.95), Inches(10), Inches(0.6),
         "appointments, insurance, and clinic info — 24 / 7.",
         size=20, color=WHITE)

# CTA pill
add_pill(s, Inches(0.8), Inches(5.9), Inches(3.2), Inches(0.45),
         "LIVE at clinic.callsphere.site", fill=TEAL, text_color=WHITE, size=12)

add_text(s, Inches(0.8), Inches(6.7), Inches(10), Inches(0.3),
         "Prepared for the Insight Healthcare team   •   May 2026",
         size=11, color=SKY)


# =========================================================
# SLIDE 2 — What it does (value, no jargon)
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, SKY)
slide_header(s, "Overview", "What the chatbot does for your patients",
             "Three simple things, done well — and it politely declines anything else.")

cards = [
    ("Book a visit",
     "Patients can see open slots and schedule, reschedule, or cancel —\nin plain English, any time of day.",
     "📅"),
    ("Check insurance",
     "Answers \"Do you take my plan?\" and looks up sample member IDs\nso patients know before they walk in.",
     "🛡️"),
    ("Answer clinic questions",
     "Hours, location, services, parking, what to bring on the first\nvisit — all from your own clinic info.",
     "💬"),
]

card_w = Inches(3.95)
card_h = Inches(3.5)
gap = Inches(0.25)
start_left = Inches(0.55)
top = Inches(2.2)

for i, (title, body, icon) in enumerate(cards):
    left = start_left + (card_w + gap) * i
    add_card(s, left, top, card_w, card_h, fill=WHITE)
    # icon circle
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                              left + Inches(0.3), top + Inches(0.3),
                              Inches(0.85), Inches(0.85))
    circ.fill.solid(); circ.fill.fore_color.rgb = MINT
    circ.line.fill.background()
    add_text(s, left + Inches(0.3), top + Inches(0.3),
             Inches(0.85), Inches(0.85), icon,
             size=28, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.3), top + Inches(1.35),
             card_w - Inches(0.6), Inches(0.5),
             title, size=20, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.95),
             card_w - Inches(0.6), Inches(1.4),
             body, size=13, color=SLATE)

# bottom note
note = add_card(s, Inches(0.55), Inches(6.0), Inches(12.25), Inches(0.7),
                fill=WHITE, border=AMBER)
add_text(s, Inches(0.85), Inches(6.05), Inches(12), Inches(0.6),
         "🔒  Privacy by design: we record which topic was asked, never the message text.",
         size=13, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 3 — How a patient uses it (simple journey)
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "Patient Journey", "How a patient uses it",
             "From question to answer in seconds — no app to download, no login.")

# 4 steps with arrows
steps = [
    ("1. Patient visits", "clinic.callsphere.site", "👤"),
    ("2. Types a question", "\"Can I book a physical\nThursday morning?\"", "⌨️"),
    ("3. AI understands", "Figures out the topic\nand finds the answer", "🧠"),
    ("4. Gets a reply", "Friendly, accurate,\nin plain English", "✅"),
]

step_w = Inches(2.7)
step_h = Inches(2.4)
arrow_w = Inches(0.45)
total_w = step_w * 4 + arrow_w * 3
start = (prs.slide_width - total_w) / 2
top = Inches(2.4)

for i, (title, body, icon) in enumerate(steps):
    left = start + (step_w + arrow_w) * i
    add_card(s, left, top, step_w, step_h, fill=SKY, border=TEAL)
    add_text(s, left, top + Inches(0.25), step_w, Inches(0.7),
             icon, size=34, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(0.95), step_w, Inches(0.4),
             title, size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.15), top + Inches(1.35),
             step_w - Inches(0.3), Inches(1.0),
             body, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    if i < 3:
        ax1 = left + step_w + Emu(20000)
        ax2 = left + step_w + arrow_w - Emu(20000)
        ay  = top + step_h / 2
        add_arrow(s, ax1, ay, ax2, ay, color=TEAL, weight=2.5)

# bottom benefit row
add_text(s, Inches(0.55), Inches(5.4), Inches(12.25), Inches(0.4),
         "Why patients like it", size=14, bold=True, color=NAVY)

benefits = [
    "Available 24 / 7 — even after the front desk closes",
    "No phone hold music, no waiting on hold",
    "Same answer every time — nothing falls through the cracks",
]
for i, b in enumerate(benefits):
    top_b = Inches(5.9) + Inches(0.45) * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(0.6), top_b + Inches(0.08),
                             Inches(0.18), Inches(0.18))
    dot.fill.solid(); dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()
    add_text(s, Inches(0.95), top_b, Inches(11.5), Inches(0.4),
             b, size=13, color=SLATE)


# =========================================================
# SLIDE 4 — Flowchart 1: Patient journey through the system
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "How it works", "Flowchart 1 — From patient to reply",
             "A peek under the hood, kept simple. Every step is on equipment we control.")

# Layout: Patient → Website → Brain (assistant) → Reply (back)
boxes = [
    ("Patient",         "asks a question",        Inches(0.7),  Inches(3.4), Inches(1.8), Inches(1.4), SKY,  TEAL),
    ("Clinic website",  "clinic.callsphere.site", Inches(3.0),  Inches(3.4), Inches(2.4), Inches(1.4), WHITE, TEAL),
    ("AI assistant",    "understands the\nquestion",       Inches(5.9), Inches(3.4), Inches(2.4), Inches(1.4), MINT, TEAL),
    ("Clinic database", "appointments,\ninsurance, FAQs",  Inches(8.8), Inches(3.4), Inches(2.4), Inches(1.4), WHITE, TEAL),
    ("Reply",           "shown to patient",       Inches(11.2), Inches(3.4), Inches(1.7), Inches(1.4), SKY, TEAL),
]

shapes_ref = []
for title, sub, l, t, w, h, fill, brd in boxes:
    shp = add_flow_box(s, l, t, w, h, title, sub, fill=fill, border=brd)
    shapes_ref.append((l, t, w, h))

# forward arrows
for i in range(len(boxes) - 1):
    l, t, w, h = shapes_ref[i]
    l2, t2, w2, h2 = shapes_ref[i + 1]
    y_mid = t + h / 2
    add_arrow(s, l + w + Emu(10000), y_mid - Inches(0.05),
              l2 - Emu(10000), y_mid - Inches(0.05),
              color=TEAL, weight=2.5)

# return path (under)
y_back = Inches(5.6)
# from last box down then back to patient
ret = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                              shapes_ref[-1][0] + shapes_ref[-1][2] / 2,
                              shapes_ref[-1][1] + shapes_ref[-1][3],
                              shapes_ref[-1][0] + shapes_ref[-1][2] / 2,
                              y_back)
ret.line.color.rgb = AMBER; ret.line.width = Pt(2)

ret2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               shapes_ref[-1][0] + shapes_ref[-1][2] / 2,
                               y_back,
                               shapes_ref[0][0] + shapes_ref[0][2] / 2,
                               y_back)
ret2.line.color.rgb = AMBER; ret2.line.width = Pt(2)

ret3 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               shapes_ref[0][0] + shapes_ref[0][2] / 2,
                               y_back,
                               shapes_ref[0][0] + shapes_ref[0][2] / 2,
                               shapes_ref[0][1] + shapes_ref[0][3] + Emu(10000))
ret3.line.color.rgb = AMBER; ret3.line.width = Pt(2)
# arrowhead on ret3
from pptx.oxml.ns import qn as _qn
ln3 = ret3.line._get_or_add_ln()
tail = ln3.makeelement(_qn('a:tailEnd'),
                        {'type': 'triangle', 'w': 'med', 'len': 'med'})
ln3.append(tail)

add_text(s, Inches(5.4), Inches(5.65), Inches(3.5), Inches(0.4),
         "answer sent back to patient", size=11, bold=True,
         color=AMBER, align=PP_ALIGN.CENTER)

# Caption strip
add_card(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.6),
         fill=SKY, border=SOFT)
add_text(s, Inches(0.85), Inches(6.3), Inches(12), Inches(0.5),
         "In plain terms: the patient's question travels to a smart helper, "
         "the helper looks up the right info, and a friendly reply comes back.",
         size=12, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 5 — Flowchart 2: How the assistant chooses an expert
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "How it works", "Flowchart 2 — Four specialists, one polite gatekeeper",
             "The AI routes every question to a focused expert — or politely declines.")

# Top: incoming question
add_flow_box(s, Inches(5.4), Inches(1.85), Inches(2.5), Inches(0.75),
             "Patient question", fill=SKY, border=TEAL,
             title_size=14)

# Triage (diamond / hex)
triage = s.shapes.add_shape(MSO_SHAPE.HEXAGON,
                             Inches(5.0), Inches(2.95),
                             Inches(3.3), Inches(1.1))
triage.fill.solid(); triage.fill.fore_color.rgb = NAVY
triage.line.fill.background()
tf = triage.text_frame
tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Smart router"
r.font.name = "Calibri"; r.font.size = Pt(15); r.font.bold = True
r.font.color.rgb = WHITE
p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
r2 = p2.add_run(); r2.text = "What is this about?"
r2.font.name = "Calibri"; r2.font.size = Pt(10)
r2.font.color.rgb = SKY

# arrow question → triage
add_arrow(s, Inches(6.65), Inches(2.62),
          Inches(6.65), Inches(2.93), color=TEAL, weight=2.5)

# Four child boxes
specialists = [
    ("Appointments",  "book, reschedule,\ncancel visits",          MINT,    Inches(0.55),  Inches(5.0)),
    ("Insurance",     "what plans you accept,\nmember coverage",   MINT,    Inches(3.85),  Inches(5.0)),
    ("Clinic info",   "hours, location,\nservices, FAQs",          MINT,    Inches(7.15),  Inches(5.0)),
    ("Politely declines", "anything off-topic\n(medical advice, etc.)", ROSE, Inches(10.45), Inches(5.0)),
]
spec_w = Inches(2.6); spec_h = Inches(1.4)

# router center for arrows
rx = Inches(6.65); ry = Inches(4.05)

for title, body, fill, l, t in specialists:
    add_flow_box(s, l, t, spec_w, spec_h, title, body,
                 fill=fill, border=TEAL, title_size=14, subtitle_size=10)
    # arrow from router to top-center of box
    tx = l + spec_w / 2
    ty = t
    add_arrow(s, rx, ry, tx, ty - Emu(5000),
              color=TEAL if "decline" not in title.lower() else AMBER,
              weight=2.0)

# legend
add_text(s, Inches(0.55), Inches(6.65), Inches(12), Inches(0.35),
         "Each specialist only answers what it knows. Out-of-scope questions get a "
         "friendly redirect, never a wrong answer.",
         size=12, color=SLATE, align=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 6 — What's safe, what's stored
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, SKY)
slide_header(s, "Trust & Safety", "What we store — and what we don't",
             "Designed for a healthcare context from day one.")

# Two columns: STORE vs DON'T STORE
col_w = Inches(5.9); col_h = Inches(4.4)
top = Inches(2.1)

# Left card (stored)
left_card = add_card(s, Inches(0.55), top, col_w, col_h, fill=WHITE, border=TEAL)
add_pill(s, Inches(0.85), top + Inches(0.25), Inches(2.0), Inches(0.4),
         "WE DO STORE", fill=TEAL, text_color=WHITE, size=11)
stored = [
    ("Clinic information", "Hours, address, services, accepted insurance"),
    ("Appointment slots & bookings", "When a visit is confirmed — name, email, reason"),
    ("Anonymous usage stats", "Which topic was asked, response speed (no message text)"),
]
for i, (h, b) in enumerate(stored):
    y = top + Inches(0.95) + Inches(1.1) * i
    add_text(s, Inches(0.95), y, col_w - Inches(0.6), Inches(0.35),
             "✓ " + h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(1.2), y + Inches(0.35), col_w - Inches(0.8),
             Inches(0.7), b, size=12, color=SLATE)

# Right card (not stored)
right_card = add_card(s, Inches(6.85), top, col_w, col_h,
                      fill=WHITE, border=AMBER)
add_pill(s, Inches(7.15), top + Inches(0.25), Inches(2.6), Inches(0.4),
         "WE DO NOT STORE", fill=AMBER, text_color=WHITE, size=11)
not_stored = [
    ("Chat message text", "The actual words a patient types are never saved"),
    ("Patient identity for browsing", "No login required — only booking captures contact info"),
    ("Raw IP addresses", "Visitor counts use a one-way hashed identifier"),
]
for i, (h, b) in enumerate(not_stored):
    y = top + Inches(0.95) + Inches(1.1) * i
    add_text(s, Inches(7.25), y, col_w - Inches(0.6), Inches(0.35),
             "✗ " + h, size=14, bold=True, color=NAVY)
    add_text(s, Inches(7.5), y + Inches(0.35), col_w - Inches(0.8),
             Inches(0.7), b, size=12, color=SLATE)

# bottom takeaway
add_card(s, Inches(0.55), Inches(6.65), Inches(12.25), Inches(0.55),
         fill=NAVY, border=NAVY)
add_text(s, Inches(0.85), Inches(6.68), Inches(12), Inches(0.5),
         "The chatbot is helpful, but it stays in its lane — and it keeps patient "
         "conversations private by default.",
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 7 — Status & what's next
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "Status", "Live today — and easy to grow",
             "Already running at clinic.callsphere.site, ready for real patient traffic.")

# Status badge
badge = add_card(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(2.2),
                 fill=MINT, border=TEAL)
add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
         "● LIVE", size=14, bold=True, color=TEAL)
add_text(s, Inches(0.85), Inches(2.55), Inches(5.5), Inches(0.6),
         "clinic.callsphere.site", size=22, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(3.15), Inches(5.5), Inches(0.9),
         "Secure (HTTPS), seeded with your clinic's hours, providers,\n"
         "accepted insurance, and 15 starter FAQs.",
         size=12, color=SLATE)

# Quick stats / scope
stats = add_card(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(2.2),
                 fill=SKY, border=TEAL)
add_text(s, Inches(7.15), Inches(2.15), Inches(5.5), Inches(0.4),
         "WHAT'S IN PRODUCTION", size=11, bold=True, color=TEAL)
items = [
    "3 family-medicine providers + 2 weeks of slots",
    "Blue Cross Blue Shield PPO sample plan",
    "15 clinic FAQs (hours, parking, telehealth, …)",
    "Friendly tone, refuses medical advice",
]
for i, it in enumerate(items):
    add_text(s, Inches(7.35), Inches(2.55) + Inches(0.36) * i,
             Inches(5.4), Inches(0.36),
             "•  " + it, size=12, color=NAVY)

# Roadmap card (next)
add_text(s, Inches(0.55), Inches(4.5), Inches(12), Inches(0.4),
         "What we can add next (when you're ready)",
         size=16, bold=True, color=NAVY)

next_steps = [
    ("More providers & specialties", "Add new doctors or service lines — slots regenerate automatically."),
    ("More insurance plans",          "Expand the accepted-payer list as your network grows."),
    ("Email / SMS confirmations",     "Send a booking confirmation to the patient and the front desk."),
    ("Admin dashboard",               "A simple page for staff to view bookings and analytics."),
]
card_w = Inches(2.95); card_h = Inches(2.0)
gap = Inches(0.15)
total = card_w * 4 + gap * 3
start = (prs.slide_width - total) / 2
top = Inches(5.0)
for i, (title, body) in enumerate(next_steps):
    left = start + (card_w + gap) * i
    add_card(s, left, top, card_w, card_h, fill=WHITE, border=SOFT)
    # number badge
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             left + Inches(0.2), top + Inches(0.2),
                             Inches(0.45), Inches(0.45))
    num.fill.solid(); num.fill.fore_color.rgb = TEAL; num.line.fill.background()
    add_text(s, left + Inches(0.2), top + Inches(0.2),
             Inches(0.45), Inches(0.45), str(i + 1),
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.75), top + Inches(0.2),
             card_w - Inches(0.95), Inches(0.5),
             title, size=13, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), top + Inches(0.85),
             card_w - Inches(0.4), card_h - Inches(1.0),
             body, size=11, color=SLATE)


# =========================================================
# SLIDE 8 — Thank you / contact
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, NAVY)

# decorative
blob = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(-3), Inches(-3), Inches(8), Inches(8))
blob.fill.solid(); blob.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x60)
blob.line.fill.background()
blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(9), Inches(4), Inches(6), Inches(6))
blob2.fill.solid(); blob2.fill.fore_color.rgb = TEAL
blob2.line.fill.background()

add_text(s, Inches(0.8), Inches(2.6), Inches(12), Inches(1.2),
         "Thank you.", size=64, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.9), Inches(12), Inches(0.6),
         "Try the live chatbot today — and tell us what you'd change.",
         size=20, color=SKY)

# CTA card
cta = add_card(s, Inches(0.8), Inches(5.0), Inches(7.0), Inches(1.4),
               fill=WHITE, border=TEAL)
add_text(s, Inches(1.0), Inches(5.15), Inches(6.5), Inches(0.4),
         "VISIT", size=11, bold=True, color=TEAL)
add_text(s, Inches(1.0), Inches(5.45), Inches(6.5), Inches(0.6),
         "clinic.callsphere.site", size=24, bold=True, color=NAVY)
add_text(s, Inches(1.0), Inches(5.95), Inches(6.5), Inches(0.4),
         "Open it on any phone or laptop — no login needed.",
         size=12, color=SLATE)


# ---------- Save ----------
out = "/home/ubuntu/apps/insight_healthcare/Insight_Healthcare_Chatbot_Overview.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
