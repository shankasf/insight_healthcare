"""Business-first deck for Insight Healthcare Chatbot.

Narrative: Problem -> Solution -> How it works (simple) -> Value -> What's next.
Audience: Non-technical clinic owner / operator.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn


# ---------- Palette ----------
NAVY    = RGBColor(0x0F, 0x2A, 0x4A)
TEAL    = RGBColor(0x1F, 0x8F, 0x8F)
SKY     = RGBColor(0xE6, 0xF2, 0xF7)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
SLATE   = RGBColor(0x4A, 0x5A, 0x6A)
SOFT    = RGBColor(0xD9, 0xE6, 0xEE)
AMBER   = RGBColor(0xE8, 0x9B, 0x3C)
MINT    = RGBColor(0xC7, 0xE9, 0xD9)
ROSE    = RGBColor(0xF7, 0xD7, 0xD7)
CORAL   = RGBColor(0xE7, 0x6F, 0x6F)
GOLD    = RGBColor(0xF6, 0xC8, 0x5F)


# ---------- Helpers ----------
def set_slide_bg(slide, color):
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0,
        prs.slide_width, prs.slide_height,
    )
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)


def add_text(slide, left, top, width, height, text,
             size=18, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Calibri", italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height,
             fill=WHITE, border=SOFT, radius=True, border_w=0.75):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        left, top, width, height,
    )
    shape.adjustments[0] = 0.06 if radius else 0
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border
    shape.line.width = Pt(border_w)
    return shape


def add_pill(slide, left, top, width, height, text,
             fill=MINT, text_color=NAVY, size=11, bold=True):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    pill.adjustments[0] = 0.5
    pill.fill.solid()
    pill.fill.fore_color.rgb = fill
    pill.line.fill.background()
    tf = pill.text_frame
    tf.margin_left = Inches(0.08); tf.margin_right = Inches(0.08)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Calibri"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = text_color
    return pill


def add_arrow(slide, x1, y1, x2, y2, color=TEAL, weight=2.0):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    ln = line.line._get_or_add_ln()
    tail = ln.find(qn('a:tailEnd'))
    if tail is None:
        tail = ln.makeelement(qn('a:tailEnd'),
                              {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    return line


def slide_header(slide, eyebrow, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  0, 0, prs.slide_width, Inches(0.18))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()

    add_text(slide, Inches(0.55), Inches(0.32), Inches(10), Inches(0.3),
             eyebrow.upper(), size=11, bold=True, color=TEAL)
    add_text(slide, Inches(0.55), Inches(0.55), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=NAVY)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(1.05), Inches(12), Inches(0.4),
                 subtitle, size=14, color=SLATE)

    add_text(slide, Inches(0.55), Inches(7.05), Inches(8), Inches(0.3),
             "Insight Healthcare  •  AI Front Desk Assistant",
             size=9, color=SLATE)
    add_text(slide, Inches(11.5), Inches(7.05), Inches(1.7), Inches(0.3),
             "clinic.callsphere.site", size=9, color=TEAL, align=PP_ALIGN.RIGHT)


def add_flow_box(slide, left, top, width, height, title,
                 subtitle=None, fill=WHITE, border=TEAL,
                 title_color=NAVY, subtitle_color=SLATE,
                 title_size=14, subtitle_size=10):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  left, top, width, height)
    box.adjustments[0] = 0.15
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = border
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08); tf.margin_bottom = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    r.font.name = "Calibri"
    r.font.size = Pt(title_size)
    r.font.bold = True
    r.font.color.rgb = title_color
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = "Calibri"
        r2.font.size = Pt(subtitle_size)
        r2.font.color.rgb = subtitle_color
    return box


# ---------- Build ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


# =========================================================
# SLIDE 1 — Title
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, NAVY)

blob = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(9.5), Inches(-2), Inches(7), Inches(7))
blob.fill.solid(); blob.fill.fore_color.rgb = TEAL
blob.line.fill.background()
blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(-2), Inches(5), Inches(5), Inches(5))
blob2.fill.solid(); blob2.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x60)
blob2.line.fill.background()

add_text(s, Inches(0.8), Inches(1.8), Inches(10), Inches(0.5),
         "INSIGHT HEALTHCARE", size=14, bold=True, color=TEAL)
add_text(s, Inches(0.8), Inches(2.35), Inches(11), Inches(2.0),
         "Fewer missed calls.\nMore booked visits.",
         size=52, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(4.95), Inches(11), Inches(0.6),
         "An AI front-desk assistant that answers patient questions",
         size=20, color=WHITE)
add_text(s, Inches(0.8), Inches(5.3), Inches(11), Inches(0.6),
         "and books appointments — 24 / 7.",
         size=20, color=WHITE)

add_pill(s, Inches(0.8), Inches(6.2), Inches(3.2), Inches(0.45),
         "LIVE at clinic.callsphere.site", fill=TEAL, text_color=WHITE, size=12)
add_text(s, Inches(0.8), Inches(6.85), Inches(10), Inches(0.3),
         "Prepared for the Insight Healthcare team   •   May 2026",
         size=11, color=SKY)


# =========================================================
# SLIDE 2 — THE PROBLEM
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, SKY)
slide_header(s, "The problem", "Patients can't get help when they need it",
             "Most patient questions arrive outside business hours — and the front desk can't be in two places at once.")

# Big stat band
band = add_card(s, Inches(0.55), Inches(1.85), Inches(12.25), Inches(1.5),
                fill=NAVY, border=NAVY)
add_text(s, Inches(0.85), Inches(1.95), Inches(11.5), Inches(0.5),
         "Today, a clinic typically loses…", size=13, color=SKY,
         anchor=MSO_ANCHOR.TOP)
add_text(s, Inches(0.85), Inches(2.4), Inches(11.5), Inches(0.9),
         "30% of inbound questions to voicemail, missed calls, "
         "and slow email replies.",
         size=22, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)

# 3 pain cards
pains = [
    ("After-hours questions go unanswered",
     "Patients ask about hours, parking, insurance, or booking at 9 PM — and reach voicemail.",
     "🌙"),
    ("Front desk is overloaded by phone",
     "Receptionists answer the same handful of questions all day instead of caring for in-person patients.",
     "📞"),
    ("Patients drop off before booking",
     "If it isn't easy to book, patients try another clinic. Each missed booking is real lost revenue.",
     "📉"),
]
card_w = Inches(3.95); card_h = Inches(3.0); gap = Inches(0.25)
start = Inches(0.55); top = Inches(3.7)
for i, (title, body, icon) in enumerate(pains):
    left = start + (card_w + gap) * i
    add_card(s, left, top, card_w, card_h, fill=WHITE, border=CORAL, border_w=1.5)
    add_text(s, left + Inches(0.3), top + Inches(0.25),
             Inches(0.7), Inches(0.7), icon, size=30,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.3), top + Inches(1.05),
             card_w - Inches(0.6), Inches(0.6),
             title, size=15, bold=True, color=NAVY)
    add_text(s, left + Inches(0.3), top + Inches(1.65),
             card_w - Inches(0.6), card_h - Inches(1.8),
             body, size=12, color=SLATE)

add_text(s, Inches(0.55), Inches(6.95), Inches(12), Inches(0.3),
         "Net effect: lower bookings, frustrated patients, and an overworked front desk.",
         size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 3 — THE SOLUTION
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "The solution", "An AI assistant on your clinic's website",
             "Patients ask questions in plain English. The assistant answers — and books visits — instantly.")

# Hero illustration card (left)
hero = add_card(s, Inches(0.55), Inches(1.95), Inches(5.6), Inches(4.8),
                fill=SKY, border=TEAL)

# fake phone/chat mock inside
phone = add_card(s, Inches(1.35), Inches(2.2), Inches(4.0), Inches(4.3),
                 fill=WHITE, border=SOFT, border_w=1.0)
add_text(s, Inches(1.55), Inches(2.35), Inches(3.6), Inches(0.35),
         "Insight Healthcare", size=11, bold=True, color=NAVY)
add_text(s, Inches(1.55), Inches(2.6), Inches(3.6), Inches(0.3),
         "Online now", size=9, color=TEAL)

# user bubble
ub = add_card(s, Inches(2.0), Inches(3.05), Inches(3.2), Inches(0.7),
              fill=MINT, border=MINT)
add_text(s, Inches(2.15), Inches(3.1), Inches(3.0), Inches(0.6),
         "\"Can I book a physical Thursday morning?\"",
         size=10, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# bot bubble
bb = add_card(s, Inches(1.55), Inches(3.95), Inches(3.4), Inches(1.1),
              fill=SKY, border=SOFT)
add_text(s, Inches(1.7), Inches(4.0), Inches(3.2), Inches(1.0),
         "Yes — Dr. Chen has 9:00 and 10:30 open Thursday. "
         "Want me to book one?",
         size=10, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# user bubble 2
ub2 = add_card(s, Inches(2.4), Inches(5.2), Inches(2.8), Inches(0.55),
              fill=MINT, border=MINT)
add_text(s, Inches(2.5), Inches(5.22), Inches(2.65), Inches(0.5),
         "\"10:30 please.\"",
         size=10, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)

# bot confirm
bb2 = add_card(s, Inches(1.55), Inches(5.9), Inches(3.4), Inches(0.55),
              fill=SKY, border=SOFT)
add_text(s, Inches(1.7), Inches(5.92), Inches(3.2), Inches(0.5),
         "Booked ✓  You'll get a confirmation email.",
         size=10, bold=True, color=TEAL, anchor=MSO_ANCHOR.MIDDLE)

# Right: 3 promises
add_text(s, Inches(6.55), Inches(1.95), Inches(6.3), Inches(0.5),
         "What patients can do, instantly:",
         size=18, bold=True, color=NAVY)

promises = [
    ("Book, reschedule, or cancel a visit",
     "Sees real open slots from your schedule and confirms in seconds."),
    ("Check if you take their insurance",
     "Answers \"Do you take BCBS PPO?\" and looks up sample member coverage."),
    ("Get clinic info anytime",
     "Hours, location, parking, services, what to bring on the first visit."),
]
for i, (h, b) in enumerate(promises):
    top_p = Inches(2.55) + Inches(1.35) * i
    # check icon
    check = s.shapes.add_shape(MSO_SHAPE.OVAL,
                               Inches(6.55), top_p + Inches(0.05),
                               Inches(0.55), Inches(0.55))
    check.fill.solid(); check.fill.fore_color.rgb = TEAL
    check.line.fill.background()
    add_text(s, Inches(6.55), top_p + Inches(0.05),
             Inches(0.55), Inches(0.55), "✓",
             size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.3), top_p, Inches(5.5), Inches(0.4),
             h, size=15, bold=True, color=NAVY)
    add_text(s, Inches(7.3), top_p + Inches(0.4),
             Inches(5.5), Inches(0.8),
             b, size=12, color=SLATE)

add_card(s, Inches(0.55), Inches(6.85), Inches(12.25), Inches(0.45),
         fill=NAVY, border=NAVY)
add_text(s, Inches(0.85), Inches(6.87), Inches(12), Inches(0.4),
         "Lives on your website — no app to download, no login required.",
         size=12, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 4 — HOW IT WORKS (simple flow)
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "How it works", "From patient question to confirmed answer",
             "Four steps. No phone calls, no waiting on hold, no human in the loop for routine questions.")

steps = [
    ("Patient asks",        "on your clinic\nwebsite",        "👤", SKY),
    ("Assistant understands",  "figures out what\nthey need",    "🧠", MINT),
    ("Finds the answer",    "from your clinic's\nschedule & info", "🔎", SKY),
    ("Replies in seconds",  "or books the\nappointment",       "✅", MINT),
]
step_w = Inches(2.7); step_h = Inches(2.6); arrow_w = Inches(0.45)
total = step_w * 4 + arrow_w * 3
start = (prs.slide_width - total) / 2
top = Inches(2.3)

for i, (title, body, icon, fill) in enumerate(steps):
    left = start + (step_w + arrow_w) * i
    add_card(s, left, top, step_w, step_h, fill=fill, border=TEAL, border_w=1.5)
    # number badge
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             left + Inches(0.15), top + Inches(0.15),
                             Inches(0.4), Inches(0.4))
    num.fill.solid(); num.fill.fore_color.rgb = NAVY
    num.line.fill.background()
    add_text(s, left + Inches(0.15), top + Inches(0.15),
             Inches(0.4), Inches(0.4), str(i+1),
             size=12, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(0.55), step_w, Inches(0.75),
             icon, size=36, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left, top + Inches(1.35), step_w, Inches(0.45),
             title, size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, left + Inches(0.15), top + Inches(1.8),
             step_w - Inches(0.3), Inches(0.7),
             body, size=11, color=SLATE, align=PP_ALIGN.CENTER)
    if i < 3:
        ay = top + step_h / 2
        add_arrow(s, left + step_w + Emu(20000), ay,
                  left + step_w + arrow_w - Emu(20000), ay,
                  color=TEAL, weight=2.5)

# bottom info card
add_card(s, Inches(0.55), Inches(5.5), Inches(12.25), Inches(1.5),
         fill=SKY, border=TEAL)
add_text(s, Inches(0.85), Inches(5.6), Inches(11.5), Inches(0.4),
         "What makes it trustworthy",
         size=14, bold=True, color=NAVY)
trust_pts = [
    "Stays in its lane — only answers about appointments, insurance, and clinic info.",
    "Politely refuses medical advice and anything off-topic.",
    "Never stores the words a patient types. Only \"which topic\" and how fast it answered.",
]
for i, t in enumerate(trust_pts):
    y = Inches(6.0) + Inches(0.32) * i
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(0.85), y + Inches(0.08),
                             Inches(0.14), Inches(0.14))
    dot.fill.solid(); dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()
    add_text(s, Inches(1.1), y, Inches(11.5), Inches(0.3),
             t, size=11, color=SLATE)


# =========================================================
# SLIDE 5 — BUSINESS VALUE
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, SKY)
slide_header(s, "The value", "What this means for your business",
             "Three concrete wins — for the clinic, the front desk, and the patient.")

values = [
    ("More bookings",
     "Capture patients who land on the site after-hours instead of losing them to voicemail.",
     "💰", GOLD, "Revenue"),
    ("Lighter front desk",
     "Receptionists stop repeating the same answers all day and focus on in-person patients.",
     "🧑‍⚕️", TEAL, "Operations"),
    ("Happier patients",
     "Instant answers and same-tone replies, 24 / 7. No phone tree, no waiting.",
     "😊", CORAL, "Experience"),
]
card_w = Inches(3.95); card_h = Inches(4.0); gap = Inches(0.25)
start = Inches(0.55); top = Inches(2.0)

for i, (title, body, icon, accent, tag) in enumerate(values):
    left = start + (card_w + gap) * i
    add_card(s, left, top, card_w, card_h, fill=WHITE, border=accent, border_w=2.0)
    # accent strip
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               left + Inches(0.15), top + Inches(0.25),
                               Inches(0.12), Inches(0.7))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background()
    add_pill(s, left + Inches(0.4), top + Inches(0.3),
             Inches(1.5), Inches(0.32),
             tag.upper(), fill=accent, text_color=WHITE, size=9)
    add_text(s, left + Inches(0.4), top + Inches(0.8),
             card_w - Inches(0.8), Inches(0.8),
             icon, size=44)
    add_text(s, left + Inches(0.4), top + Inches(1.85),
             card_w - Inches(0.8), Inches(0.55),
             title, size=20, bold=True, color=NAVY)
    add_text(s, left + Inches(0.4), top + Inches(2.45),
             card_w - Inches(0.8), Inches(1.5),
             body, size=12, color=SLATE)

# bottom strip — what doesn't change
strip = add_card(s, Inches(0.55), Inches(6.25), Inches(12.25), Inches(0.7),
                 fill=NAVY, border=NAVY)
add_text(s, Inches(0.85), Inches(6.3), Inches(12), Inches(0.6),
         "And nothing risky changes: no medical advice, no chat content stored, "
         "and a human can always take over.",
         size=13, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# =========================================================
# SLIDE 6 — VISUAL: BEFORE vs AFTER
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "Before & after", "What changes once it's live",
             "Same clinic. Same staff. A very different patient experience.")

# Before column
before = add_card(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(4.7),
                  fill=WHITE, border=CORAL, border_w=1.5)
add_pill(s, Inches(0.85), Inches(2.2), Inches(1.6), Inches(0.4),
         "BEFORE", fill=CORAL, text_color=WHITE, size=11)
add_text(s, Inches(0.85), Inches(2.75), Inches(5.5), Inches(0.5),
         "Phone-only front desk", size=18, bold=True, color=NAVY)

before_pts = [
    ("⏰", "Open only 9–5, Mon–Fri"),
    ("📞", "Patients wait on hold or leave voicemail"),
    ("🔁", "Staff answer the same questions all day"),
    ("📉", "After-hours bookings = lost bookings"),
    ("📝", "Manual scheduling in a paper / spreadsheet"),
]
for i, (icon, txt) in enumerate(before_pts):
    y = Inches(3.4) + Inches(0.55) * i
    add_text(s, Inches(0.95), y, Inches(0.5), Inches(0.4),
             icon, size=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(1.5), y, Inches(4.9), Inches(0.4),
             txt, size=13, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

# Arrow between
arr_mid_x = Inches(6.65)
arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                          Inches(6.45), Inches(4.0),
                          Inches(0.65), Inches(0.55))
arr.fill.solid(); arr.fill.fore_color.rgb = TEAL
arr.line.fill.background()

# After column
after = add_card(s, Inches(7.0), Inches(2.0), Inches(5.8), Inches(4.7),
                 fill=WHITE, border=TEAL, border_w=1.5)
add_pill(s, Inches(7.3), Inches(2.2), Inches(1.6), Inches(0.4),
         "AFTER", fill=TEAL, text_color=WHITE, size=11)
add_text(s, Inches(7.3), Inches(2.75), Inches(5.5), Inches(0.5),
         "Phone + AI assistant", size=18, bold=True, color=NAVY)

after_pts = [
    ("🌐", "Open 24 / 7 on your website"),
    ("⚡", "Patients get answers in seconds"),
    ("🧑‍⚕️", "Staff freed up for in-person care"),
    ("📈", "After-hours visits get booked, not lost"),
    ("📅", "Real, live schedule — no double-booking"),
]
for i, (icon, txt) in enumerate(after_pts):
    y = Inches(3.4) + Inches(0.55) * i
    add_text(s, Inches(7.4), y, Inches(0.5), Inches(0.4),
             icon, size=18, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, Inches(7.95), y, Inches(4.7), Inches(0.4),
             txt, size=13, color=SLATE, anchor=MSO_ANCHOR.MIDDLE)

add_text(s, Inches(0.55), Inches(6.85), Inches(12), Inches(0.3),
         "The front desk doesn't go away — it just stops being the only door.",
         size=12, italic=True, color=SLATE, align=PP_ALIGN.CENTER)


# =========================================================
# SLIDE 7 — STATUS + WHAT'S NEXT
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, WHITE)
slide_header(s, "Status", "Live today — and easy to grow",
             "Running at clinic.callsphere.site, seeded with your clinic's real info.")

# Status badge
add_card(s, Inches(0.55), Inches(2.0), Inches(6.0), Inches(2.2),
        fill=MINT, border=TEAL, border_w=1.5)
add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.4),
         "● LIVE", size=14, bold=True, color=TEAL)
add_text(s, Inches(0.85), Inches(2.55), Inches(5.5), Inches(0.6),
         "clinic.callsphere.site", size=22, bold=True, color=NAVY)
add_text(s, Inches(0.85), Inches(3.15), Inches(5.5), Inches(0.9),
         "Secure (HTTPS), seeded with your clinic's hours, providers,\n"
         "accepted insurance, and 15 starter FAQs.",
         size=12, color=SLATE)

# Production scope
add_card(s, Inches(6.85), Inches(2.0), Inches(6.0), Inches(2.2),
         fill=SKY, border=TEAL, border_w=1.5)
add_text(s, Inches(7.15), Inches(2.15), Inches(5.5), Inches(0.4),
         "WHAT'S IN PRODUCTION", size=11, bold=True, color=TEAL)
items = [
    "3 family-medicine providers + 2 weeks of slots",
    "Blue Cross Blue Shield PPO sample plan",
    "15 clinic FAQs (hours, parking, telehealth, …)",
    "Friendly tone, refuses medical advice",
]
for i, it in enumerate(items):
    add_text(s, Inches(7.35), Inches(2.55) + Inches(0.36) * i,
             Inches(5.4), Inches(0.36),
             "•  " + it, size=12, color=NAVY)

# What's next
add_text(s, Inches(0.55), Inches(4.5), Inches(12), Inches(0.4),
         "When you're ready, we can add", size=16, bold=True, color=NAVY)

next_steps = [
    ("More providers & specialties",
     "Add new doctors or service lines — slots regenerate automatically."),
    ("More insurance plans",
     "Expand the accepted-payer list as your network grows."),
    ("Email / SMS confirmations",
     "Booking confirmations sent to the patient and the front desk."),
    ("Staff dashboard",
     "Simple admin view for bookings, demand, and trends."),
]
card_w = Inches(2.95); card_h = Inches(2.0); gap = Inches(0.15)
total = card_w * 4 + gap * 3
start = (prs.slide_width - total) / 2
top = Inches(5.0)
for i, (title, body) in enumerate(next_steps):
    left = start + (card_w + gap) * i
    add_card(s, left, top, card_w, card_h, fill=WHITE, border=SOFT)
    num = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             left + Inches(0.2), top + Inches(0.2),
                             Inches(0.45), Inches(0.45))
    num.fill.solid(); num.fill.fore_color.rgb = TEAL; num.line.fill.background()
    add_text(s, left + Inches(0.2), top + Inches(0.2),
             Inches(0.45), Inches(0.45), str(i + 1),
             size=14, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.75), top + Inches(0.2),
             card_w - Inches(0.95), Inches(0.5),
             title, size=13, bold=True, color=NAVY,
             anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, left + Inches(0.2), top + Inches(0.85),
             card_w - Inches(0.4), card_h - Inches(1.0),
             body, size=11, color=SLATE)


# =========================================================
# SLIDE 8 — CLOSING / CALL TO ACTION
# =========================================================
s = prs.slides.add_slide(blank)
set_slide_bg(s, NAVY)

blob = s.shapes.add_shape(MSO_SHAPE.OVAL,
                          Inches(-3), Inches(-3), Inches(8), Inches(8))
blob.fill.solid(); blob.fill.fore_color.rgb = RGBColor(0x16, 0x3A, 0x60)
blob.line.fill.background()
blob2 = s.shapes.add_shape(MSO_SHAPE.OVAL,
                           Inches(9), Inches(4), Inches(6), Inches(6))
blob2.fill.solid(); blob2.fill.fore_color.rgb = TEAL
blob2.line.fill.background()

add_text(s, Inches(0.8), Inches(2.0), Inches(12), Inches(1.0),
         "Ready when you are.", size=54, bold=True, color=WHITE)
add_text(s, Inches(0.8), Inches(3.25), Inches(12), Inches(0.6),
         "The chatbot is live. The next step is yours —",
         size=20, color=SKY)
add_text(s, Inches(0.8), Inches(3.65), Inches(12), Inches(0.6),
         "try it, tell us what to refine, and we'll iterate.",
         size=20, color=SKY)

cta = add_card(s, Inches(0.8), Inches(4.8), Inches(7.0), Inches(1.5),
               fill=WHITE, border=TEAL)
add_text(s, Inches(1.0), Inches(4.95), Inches(6.5), Inches(0.4),
         "TRY IT NOW", size=11, bold=True, color=TEAL)
add_text(s, Inches(1.0), Inches(5.25), Inches(6.5), Inches(0.6),
         "clinic.callsphere.site", size=24, bold=True, color=NAVY)
add_text(s, Inches(1.0), Inches(5.8), Inches(6.5), Inches(0.4),
         "Open it on any phone or laptop — no login needed.",
         size=12, color=SLATE)

add_text(s, Inches(0.8), Inches(6.7), Inches(12), Inches(0.35),
         "Questions? sagar@callsphere.tech",
         size=12, color=SKY)


# ---------- Save ----------
out = "/home/ubuntu/apps/insight_healthcare/Insight_Healthcare_Business_Pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
